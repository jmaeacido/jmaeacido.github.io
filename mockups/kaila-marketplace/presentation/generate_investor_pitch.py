from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from PIL import Image, ImageDraw, ImageFont, JpegImagePlugin
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
APP_ROOT = ROOT.parent
PACKAGE_DIR = ROOT / "investor-package"
PPTX_PATH = ROOT / "KAILA_Investor_Pitch_Deck.pptx"
PDF_PATH = ROOT / "KAILA_Investor_Pitch_Deck.pdf"
XLSX_PATH = PACKAGE_DIR / "KAILA_Financial_Model_Starter.xlsx"

LOGO = APP_ROOT / "assets" / "kaila-preview.png"
ICON = APP_ROOT / "assets" / "android-chrome-512x512.png"

W, H = 13.333, 7.5

COLORS = {
    "ink": "10242A",
    "muted": "567179",
    "teal": "0B4552",
    "sea": "0E6D72",
    "mint": "1FAE9B",
    "gold": "F4B942",
    "coral": "E86F51",
    "cream": "F7F3EA",
    "paper": "FBFAF7",
    "line": "D7E3E1",
    "white": "FFFFFF",
}

FOUNDER = "John Mark Agustin E. Acido"
EMAIL = "94jmaea94@gmail.com"
PILOT_CITY = "Gingoog City"


SLIDES = [
    {
        "type": "title",
        "kicker": "Investor Pitch",
        "title": "KAILA",
        "subtitle": "The On-Demand Services Marketplace for the Philippines",
        "tagline": "Kung wala kay kaila, naa mi.",
        "details": [FOUNDER, EMAIL, f"Pilot City: {PILOT_CITY}"],
    },
    {
        "type": "split",
        "kicker": "Problem",
        "title": "Finding trusted local help is still informal, slow, and hard to verify.",
        "left_title": "Customers",
        "left": [
            "Hard to find trusted workers quickly",
            "No standardized pricing or offer comparison",
            "No clear guarantees, ratings, or work history",
            "Communication scattered across referrals, Facebook, and Messenger",
        ],
        "right_title": "Service Providers",
        "right": [
            "Depend heavily on referrals",
            "Inconsistent job flow and income",
            "Limited digital presence",
            "Difficult to acquire customers outside their network",
        ],
    },
    {
        "type": "flow",
        "kicker": "Solution",
        "title": "KAILA turns informal referrals into a structured request, match, pay, and review workflow.",
        "steps": [
            ("1", "Request", "Customer opens KAILA and posts the service need."),
            ("2", "Match", "Relevant local providers receive the request."),
            ("3", "Work", "Provider performs the agreed service."),
            ("4", "Payment", "Payment is completed or recorded."),
            ("5", "Review", "Both sides build visible reputation."),
        ],
    },
    {
        "type": "market",
        "kicker": "Market Opportunity",
        "title": "Philippines is the upside market; Gingoog City is the proof market.",
        "metrics": [
            ("TAM", "Philippine local services spend", "Est. PHP 250B+ annual addressable household and small-business service activity."),
            ("SAM", "Digitally reachable early categories", "Est. PHP 25B-50B across repair, maintenance, cleaning, digital services, and skilled work."),
            ("SOM", "3-year realistic capture", "Est. PHP 20M-60M GMV from disciplined city-by-city expansion after Gingoog validation."),
        ],
        "note": "Sizing is an explicit planning estimate using PSA household/service spend categories, DTI MSME data, and KAILA category assumptions.",
    },
    {
        "type": "demo",
        "kicker": "Product Demo",
        "title": "Early mockups show the core marketplace workflow before full app build.",
        "bullets": [
            "Request board for client needs and provider offers",
            "Provider directory with categories, areas, and trust signals",
            "Admin-style dashboard for pilot operations and manual matching",
            "Mobile-first experience planned for clients and providers",
        ],
    },
    {
        "type": "table",
        "kicker": "Business Model",
        "title": "KAILA monetizes after the marketplace proves real provider value.",
        "headers": ["Revenue Stream", "Description"],
        "rows": [
            ("Service fee", "% commission per completed transaction after trust is proven"),
            ("Featured listings", "Provider promotion and profile boosts"),
            ("Subscription", "Premium provider accounts with better visibility"),
            ("Ads", "Local business promotions and category sponsorships"),
        ],
    },
    {
        "type": "cards",
        "kicker": "Why Now",
        "title": "The behavior is already digital; the trust system is missing.",
        "cards": [
            ("97.5M internet users", "DataReportal reported 83.8% internet penetration in the Philippines in January 2025."),
            ("142M mobile connections", "Mobile connectivity is larger than the population, making mobile-first matching practical."),
            ("57.4% digital payment volume", "BSP reported digital retail payments reached 57.4% of transaction volume in 2024."),
            ("MSME-heavy economy", "DTI reports MSMEs make up 99.66% of Philippine business establishments."),
        ],
    },
    {
        "type": "matrix",
        "kicker": "Competition",
        "title": "KAILA wins by combining local focus, search, verification, ratings, and structured workflows.",
        "headers": ["Feature", "Facebook", "Generic Marketplace", "KAILA"],
        "rows": [
            ("Searchability", "No", "Yes", "Yes"),
            ("Ratings", "No", "Yes", "Yes"),
            ("Verification", "No", "Partial", "Yes"),
            ("Local focus", "Partial", "Partial", "Yes"),
            ("Request-to-offer workflow", "No", "Partial", "Yes"),
            ("Pilot operations support", "No", "No", "Yes"),
        ],
    },
    {
        "type": "gtm",
        "kicker": "Go-To-Market",
        "title": "Start narrow in Gingoog City, prove density, then replicate the playbook.",
        "demand": ["Facebook groups", "Barangay partnerships", "Local ads and referrals", "Small business outreach"],
        "supply": ["Skilled workers", "Freelancers", "Tradespeople", "Repair shops and service microbusinesses"],
        "pilot": ["Recruit providers first", "Run concierge matching", "Measure response and completion", "Convert proof into MVP scope"],
    },
    {
        "type": "traction",
        "kicker": "Traction",
        "title": "Current stage: founder-grade package, prototype mockups, and Gingoog validation plan.",
        "metrics": [
            ("50", "Target providers"),
            ("100", "Client survey responses"),
            ("30", "Real service requests"),
            ("15+", "Completed or strongly matched jobs"),
            ("5-8", "Starting service categories"),
            ("30-45", "Active pilot days"),
        ],
        "note": "If actual interviews, waitlist, survey results, or pilot jobs are added, this becomes the strongest investor proof slide.",
    },
    {
        "type": "financial",
        "kicker": "Financial Projections",
        "title": "3-year model: realistic ramp from concierge pilot to paid marketplace.",
        "headers": ["Metric", "Year 1", "Year 2", "Year 3"],
        "rows": [
            ("Active providers", "150", "600", "1,800"),
            ("Completed jobs", "900", "7,200", "32,400"),
            ("GMV", "PHP 0.9M", "PHP 8.6M", "PHP 48.6M"),
            ("Revenue", "PHP 0.2M", "PHP 1.7M", "PHP 7.3M"),
            ("Expenses", "PHP 1.8M", "PHP 3.6M", "PHP 7.2M"),
            ("Profit / Loss", "(PHP 1.6M)", "(PHP 1.9M)", "PHP 0.1M"),
        ],
        "note": "Assumptions are included in the investor package workbook and should be updated after the Gingoog pilot.",
    },
    {
        "type": "funding",
        "kicker": "Funding Ask",
        "title": "Raising PHP 2,000,000 to validate Gingoog City and build the first focused MVP.",
        "metrics": [
            ("40%", "Development", "PHP 800,000"),
            ("35%", "Marketing", "PHP 700,000"),
            ("15%", "Operations", "PHP 300,000"),
            ("10%", "Legal/Admin", "PHP 200,000"),
        ],
        "note": "Scenario options in the package show lean, base, and accelerated funding paths.",
    },
    {
        "type": "team",
        "kicker": "Team",
        "title": "Founder-led by a technical builder with validation-first discipline.",
        "name": FOUNDER,
        "roles": [
            "Product and technology lead",
            "Web development and prototype execution",
            "Founder-grade business documentation and MVP planning",
            "Local pilot design for Gingoog City",
        ],
        "open_roles": ["Operations lead", "Provider acquisition lead", "Marketing/community lead"],
    },
    {
        "type": "vision",
        "kicker": "Vision",
        "title": "Become the leading on-demand services platform for Philippine cities.",
        "bullets": [
            "Start with one disciplined pilot in Gingoog City.",
            "Prove matching, trust, ratings, and provider value manually.",
            "Build the MVP around the proven workflow.",
            "Expand city-by-city with a repeatable local operations playbook.",
        ],
    },
    {
        "type": "closing",
        "kicker": "Investor Package",
        "title": "What investors can review next",
        "items": [
            "Investor pitch deck",
            "Market sizing and assumptions memo",
            "Starter 3-year financial model",
            "Business plan outline",
            "Due diligence data room checklist",
        ],
        "note": "The immediate next milestone is the Gingoog City validation sprint.",
    },
]


def rgb(hex_value):
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def add_shape(slide, shape_type, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, color="ink", bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS[color])
        p.space_after = Pt(8)
    return box


def decorate(slide, index):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H, COLORS["paper"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 0.14, H, COLORS["teal"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 7.15, W, 0.35, COLORS["cream"])
    add_text(slide, "KAILA Investor Pitch Package", 0.45, 7.22, 3.2, 0.18, 8, "muted", bold=True)
    add_text(slide, f"{index:02d}", 12.35, 7.22, 0.45, 0.18, 8, "muted", align=PP_ALIGN.RIGHT)


def header(slide, data):
    add_text(slide, data["kicker"].upper(), 0.65, 0.42, 2.7, 0.22, 9, "sea", bold=True)
    add_text(slide, data["title"], 0.65, 0.8, 11.4, 0.85, 23, "ink", bold=True)


def add_card(slide, x, y, w, h, title, body, stripe="mint"):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, COLORS["white"], COLORS["line"])
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, 0.12, COLORS[stripe])
    add_text(slide, title, x + 0.22, y + 0.28, w - 0.45, 0.35, 15, "teal", bold=True)
    add_text(slide, body, x + 0.22, y + 0.85, w - 0.45, h - 1.1, 11.5, "ink")


def add_table(slide, headers, rows, x, y, widths, row_h=0.5, font_size=10.5):
    for c, header in enumerate(headers):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + sum(widths[:c]), y, widths[c], row_h, COLORS["teal"], COLORS["teal"])
        add_text(slide, header, x + sum(widths[:c]) + 0.08, y + 0.12, widths[c] - 0.16, 0.18, font_size, "white", bold=True)
    for r, row in enumerate(rows):
        fill = COLORS["white"] if r % 2 == 0 else COLORS["cream"]
        for c, val in enumerate(row):
            xx = x + sum(widths[:c])
            yy = y + row_h * (r + 1)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, xx, yy, widths[c], row_h, fill, COLORS["line"])
            add_text(slide, val, xx + 0.08, yy + 0.12, widths[c] - 0.16, 0.2, font_size, "ink", bold=(c == len(row) - 1))


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if data["type"] == "title":
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H, COLORS["teal"])
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 8.65, 0, 4.7, H, COLORS["cream"])
            if ICON.exists():
                slide.shapes.add_picture(str(ICON), Inches(9.55), Inches(0.85), Inches(2.2), Inches(2.2))
            add_text(slide, data["kicker"].upper(), 0.78, 0.72, 2.4, 0.25, 10, "gold", bold=True)
            add_text(slide, data["title"], 0.72, 1.35, 5.6, 0.85, 58, "white", bold=True)
            add_text(slide, data["subtitle"], 0.78, 2.45, 6.6, 0.6, 24, "white", bold=True)
            add_text(slide, data["tagline"], 0.78, 3.35, 4.3, 0.32, 16, "gold", bold=True)
            y = 5.75
            for detail in data["details"]:
                add_text(slide, detail, 0.78, y, 5.4, 0.24, 12, "white")
                y += 0.34
            add_text(slide, "Pilot proof in Gingoog City. National expansion story for the Philippines.", 9.15, 4.2, 3.0, 1.2, 24, "ink", bold=True)
            continue

        decorate(slide, idx)
        header(slide, data)

        if data["type"] == "split":
            add_card(slide, 0.85, 2.0, 5.65, 3.9, data["left_title"], "\n".join(data["left"]), "gold")
            add_card(slide, 6.85, 2.0, 5.65, 3.9, data["right_title"], "\n".join(data["right"]), "coral")

        elif data["type"] == "flow":
            for i, (num, title, body) in enumerate(data["steps"]):
                x = 0.55 + i * 2.52
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 2.25, 2.1, 2.55, COLORS["white"], COLORS["line"])
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.24, 2.56, 0.48, 0.48, COLORS["gold"])
                add_text(slide, num, x + 0.38, 2.67, 0.2, 0.18, 11, "ink", bold=True, align=PP_ALIGN.CENTER)
                add_text(slide, title, x + 0.24, 3.28, 1.55, 0.32, 15, "teal", bold=True)
                add_text(slide, body, x + 0.24, 3.88, 1.55, 0.65, 10.5, "ink")
                if i < 4:
                    add_text(slide, "->", x + 2.16, 3.32, 0.28, 0.25, 14, "muted", bold=True)

        elif data["type"] == "market":
            for i, (label, title, body) in enumerate(data["metrics"]):
                x = 0.8 + i * 4.1
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 2.0, 3.55, 3.25, COLORS["white"], COLORS["line"])
                add_text(slide, label, x + 0.24, 2.28, 1.0, 0.38, 24, ["teal", "sea", "coral"][i], bold=True)
                add_text(slide, title, x + 0.24, 2.95, 2.7, 0.45, 14, "ink", bold=True)
                add_text(slide, body, x + 0.24, 3.65, 2.75, 1.0, 11.5, "muted")
            add_text(slide, data["note"], 0.9, 6.0, 10.8, 0.5, 11, "muted")

        elif data["type"] == "demo":
            if LOGO.exists():
                slide.shapes.add_picture(str(LOGO), Inches(0.85), Inches(2.0), Inches(6.25))
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.6, 2.0, 4.5, 3.75, COLORS["white"], COLORS["line"])
            add_bullets(slide, data["bullets"], 7.95, 2.45, 3.75, 2.8, 13)

        elif data["type"] == "table":
            add_table(slide, data["headers"], data["rows"], 1.0, 2.1, [3.2, 7.8], 0.62, 12)

        elif data["type"] == "cards":
            for i, (title, body) in enumerate(data["cards"]):
                x = 0.85 + (i % 2) * 5.85
                y = 2.0 + (i // 2) * 1.85
                add_card(slide, x, y, 5.1, 1.35, title, body, ["mint", "gold", "sea", "coral"][i])

        elif data["type"] == "matrix":
            add_table(slide, data["headers"], data["rows"], 0.95, 2.0, [3.0, 2.4, 2.8, 2.4], 0.5, 10.5)

        elif data["type"] == "gtm":
            add_card(slide, 0.8, 2.0, 3.6, 3.65, "Demand Side", "\n".join(data["demand"]), "gold")
            add_card(slide, 4.85, 2.0, 3.6, 3.65, "Supply Side", "\n".join(data["supply"]), "mint")
            add_card(slide, 8.9, 2.0, 3.6, 3.65, "Gingoog Pilot", "\n".join(data["pilot"]), "coral")

        elif data["type"] == "traction":
            for i, (big, label) in enumerate(data["metrics"]):
                row, col = divmod(i, 3)
                x = 0.9 + col * 3.85
                y = 2.05 + row * 1.35
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.25, 0.95, COLORS["white"], COLORS["line"])
                add_text(slide, big, x + 0.22, y + 0.22, 0.8, 0.28, 19, "teal", bold=True)
                add_text(slide, label, x + 1.05, y + 0.2, 1.85, 0.35, 10.5, "ink", bold=True)
            add_text(slide, data["note"], 0.9, 5.55, 10.9, 0.55, 11, "muted")

        elif data["type"] == "financial":
            add_table(slide, data["headers"], data["rows"], 0.75, 2.0, [3.0, 2.4, 2.4, 2.4], 0.48, 10)
            add_text(slide, data["note"], 0.85, 5.95, 10.6, 0.5, 10.5, "muted")

        elif data["type"] == "funding":
            for i, (pct, label, amount) in enumerate(data["metrics"]):
                x = 0.9 + (i % 2) * 5.75
                y = 2.05 + (i // 2) * 1.55
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 4.95, 1.05, COLORS["white"], COLORS["line"])
                add_text(slide, pct, x + 0.24, y + 0.22, 0.8, 0.3, 22, "teal", bold=True)
                add_text(slide, label, x + 1.35, y + 0.18, 1.9, 0.24, 13, "ink", bold=True)
                add_text(slide, amount, x + 1.35, y + 0.55, 1.8, 0.22, 11, "muted")
            add_text(slide, data["note"], 0.9, 5.7, 10.8, 0.45, 11, "muted")

        elif data["type"] == "team":
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 2.05, 5.25, 3.55, COLORS["white"], COLORS["line"])
            add_text(slide, data["name"], 1.25, 2.45, 4.35, 0.35, 20, "teal", bold=True)
            add_bullets(slide, data["roles"], 1.25, 3.05, 4.3, 1.8, 12.5)
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.0, 2.05, 4.65, 3.55, COLORS["teal"])
            add_text(slide, "Next strategic hires / co-founders", 7.35, 2.48, 3.5, 0.45, 18, "white", bold=True)
            add_bullets(slide, data["open_roles"], 7.35, 3.25, 3.4, 1.25, 13, "white")

        elif data["type"] == "vision":
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 2.0, 11.1, 3.9, COLORS["teal"])
            add_bullets(slide, data["bullets"], 1.35, 2.65, 9.9, 2.35, 18, "white")

        elif data["type"] == "closing":
            for i, item in enumerate(data["items"]):
                x = 0.9 + (i % 3) * 3.85
                y = 2.0 + (i // 3) * 1.35
                add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.25, 0.95, COLORS["white"], COLORS["line"])
                add_text(slide, item, x + 0.25, y + 0.28, 2.55, 0.26, 13, "ink", bold=True)
            add_text(slide, data["note"], 0.9, 5.55, 10.8, 0.5, 13, "teal", bold=True)

    prs.save(PPTX_PATH)


def font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def build_pdf_fallback():
    pages = []
    scale = 144
    width, height = int(W * scale), int(H * scale)
    c = {name: f"#{value}" for name, value in COLORS.items()}
    for idx, data in enumerate(SLIDES, start=1):
        img = Image.new("RGB", (width, height), c["paper"])
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, 22, height), fill=c["teal"])
        draw.rectangle((0, int(height * 0.953), width, height), fill=c["cream"])
        draw.text((95, 65), data["kicker"].upper(), font=font(18, True), fill=c["sea"])
        draw.text((95, 125), data["title"][:90], font=font(42, True), fill=c["ink"])
        draw.text((65, 1040), "KAILA Investor Pitch Package", font=font(14, True), fill=c["muted"])
        draw.text((1780, 1040), f"{idx:02d}", font=font(14), fill=c["muted"])
        pages.append(img)
    pages[0].save(PDF_PATH, save_all=True, append_images=pages[1:], resolution=144)


def write_markdown(path, title, body):
    path.write_text(f"# {title}\n\n{body.strip()}\n", encoding="utf-8")


def build_docs():
    PACKAGE_DIR.mkdir(exist_ok=True)
    write_markdown(
        PACKAGE_DIR / "README.md",
        "KAILA Investor Pitch Package",
        """
This folder supports the investor pitch deck.

Core sequence for investors:

1. Pitch Deck: is this interesting?
2. Business Plan: can this become a real company?
3. Due Diligence Folder: should we invest?

Pilot city: Gingoog City.

Included files:

- `KAILA_Financial_Model_Starter.xlsx`
- `Market_Sizing_and_Assumptions.md`
- `Business_Plan_Outline.md`
- `Due_Diligence_Data_Room_Checklist.md`
- `Funding_Scenarios.md`
""",
    )
    write_markdown(
        PACKAGE_DIR / "Market_Sizing_and_Assumptions.md",
        "Market Sizing and Assumptions",
        """
KAILA should present the Philippines as the long-term opportunity and Gingoog City as the proof market.

TAM: Philippine local services spend. Planning estimate: PHP 250B+ annual addressable household and small-business service activity.

SAM: Digitally reachable early categories. Planning estimate: PHP 25B-50B across appliance repair, plumbing, electrical, computer repair, home maintenance, cleaning, digital services, and similar skilled local work.

SOM: Realistic 3-year capture. Planning estimate: PHP 20M-60M GMV if Gingoog validates and the playbook expands to more local markets.

Pilot assumptions:

- Pilot city: Gingoog City
- PSA PSGC population reference: 138,895 people in the 2024 POPCEN
- PSA PSGC structure reference: 79 barangays as of 31 July 2025
- Active matching period: 30-45 days
- Starting provider target: 50
- Client survey target: 100 responses
- Real request target: 30
- Completed or strongly matched jobs: 15+
- Starting categories: 5-8

Source anchors:

- PSA PSGC City of Gingoog: https://psa.gov.ph/classification/psgc/barangays/1004308000
- DataReportal Digital 2025 Philippines: https://datareportal.com/reports/digital-2025-philippines
- Bangko Sentral ng Pilipinas 2024 e-payments report: https://www.bsp.gov.ph/PaymentAndSettlement/2024_Report_on_E-payments_Measurement.pdf
- Department of Trade and Industry 2024 MSME statistics: https://www.dti.gov.ph/dti-knowledge-hub/dti-statistics/dti-msme-statistics/

Source facts used:

- DataReportal: 97.5M internet users, 83.8% internet penetration, 142M mobile connections in the Philippines in early 2025.
- BSP: digital retail payments reached 57.4% of total transaction volume and 59.0% of value in 2024.
- DTI: MSMEs comprise 99.66% of Philippine business establishments and contribute 66.58% to employment.
- PSA PSGC: Gingoog City had 138,895 people in the 2024 POPCEN and 79 barangays as of 31 July 2025.

Important note: KAILA's exact serviceable market should be recalculated after the Gingoog City pilot using real request frequency, category demand, completed job values, repeat usage, and willingness-to-pay.
""",
    )
    write_markdown(
        PACKAGE_DIR / "Business_Plan_Outline.md",
        "Business Plan Outline",
        """
1. Executive Summary
2. Company Overview
3. Problem and Customer Pain
4. KAILA Solution
5. Pilot City: Gingoog City
6. Market Research and Market Sizing
7. Competitive Analysis
8. Business Model
9. Go-To-Market Strategy
10. Operations Plan
11. Product and MVP Roadmap
12. Trust, Verification, Safety, and Dispute Handling
13. Marketing Plan
14. Financial Model
15. Funding Ask and Use of Funds
16. SWOT Analysis
17. Risk Analysis and Mitigation
18. Milestones
19. Team
20. Appendices
""",
    )
    write_markdown(
        PACKAGE_DIR / "Due_Diligence_Data_Room_Checklist.md",
        "Due Diligence Data Room Checklist",
        """
Legal:

- SEC registration once available
- Articles of incorporation once available
- Founder agreements
- IP assignment agreements
- Privacy policy and terms draft

Financial:

- Financial model
- Budget projections
- Funding scenarios
- Expense receipts and cap table once available

Product:

- Prototype screenshots
- User journeys
- MVP specification
- Architecture diagram
- Product roadmap

Research:

- Survey results
- Customer interviews
- Provider interviews
- Competitive analysis
- Pilot request tracker

Operations:

- Provider onboarding form
- Provider rules agreement
- Request tracker
- Job log
- Dispute log
- Weekly pilot reports
""",
    )
    write_markdown(
        PACKAGE_DIR / "Funding_Scenarios.md",
        "Funding Scenarios",
        """
Lean validation: PHP 500,000

- Manual Gingoog City pilot
- Provider recruitment and verification
- Survey work and field operations
- Low-code/manual MVP support

Base raise: PHP 2,000,000

- 40% development
- 35% marketing
- 15% operations
- 10% legal/admin

Accelerated pilot-to-MVP: PHP 5,000,000

- Full MVP build
- Stronger local marketing
- Dedicated operations support
- More categories and adjacent-area testing

Recommended current ask: PHP 2,000,000, subject to investor conversation and founder dilution targets.
""",
    )


def build_financial_model():
    PACKAGE_DIR.mkdir(exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "Assumptions"
    data = [
        ("Pilot city", PILOT_CITY),
        ("Raise target", 2000000),
        ("Starting active providers", 50),
        ("Year 1 active providers", 150),
        ("Year 2 active providers", 600),
        ("Year 3 active providers", 1800),
        ("Average job value Year 1", 1000),
        ("Average job value Year 2", 1200),
        ("Average job value Year 3", 1500),
        ("Completed jobs Year 1", 900),
        ("Completed jobs Year 2", 7200),
        ("Completed jobs Year 3", 32400),
        ("Take rate Year 1", 0.10),
        ("Take rate Year 2", 0.12),
        ("Take rate Year 3", 0.15),
    ]
    ws.append(["Assumption", "Value"])
    for row in data:
        ws.append(row)

    summary = wb.create_sheet("3-Year Summary")
    summary.append(["Metric", "Year 1", "Year 2", "Year 3"])
    summary_rows = [
        ("Active providers", 150, 600, 1800),
        ("Completed jobs", 900, 7200, 32400),
        ("GMV", 900000, 8640000, 48600000),
        ("Revenue", 180000, 1728000, 7290000),
        ("Expenses", 1800000, 3600000, 7200000),
        ("Profit / Loss", -1620000, -1872000, 90000),
    ]
    for row in summary_rows:
        summary.append(row)

    use = wb.create_sheet("Use of Funds")
    use.append(["Category", "Percent", "Amount"])
    for row in [
        ("Development", 0.40, 800000),
        ("Marketing", 0.35, 700000),
        ("Operations", 0.15, 300000),
        ("Legal/Admin", 0.10, 200000),
    ]:
        use.append(row)

    for sheet in wb.worksheets:
        for cell in sheet[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor=COLORS["teal"])
            cell.alignment = Alignment(horizontal="center")
        for col in sheet.columns:
            sheet.column_dimensions[col[0].column_letter].width = 24
        for row in sheet.iter_rows():
            for cell in row:
                cell.border = Border(left=Side(style="thin", color="D7E3E1"), right=Side(style="thin", color="D7E3E1"), top=Side(style="thin", color="D7E3E1"), bottom=Side(style="thin", color="D7E3E1"))
    wb.save(XLSX_PATH)


if __name__ == "__main__":
    build_pptx()
    build_pdf_fallback()
    build_docs()
    build_financial_model()
    print(PPTX_PATH)
    print(PDF_PATH)
    print(PACKAGE_DIR)
